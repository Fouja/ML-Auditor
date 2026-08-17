"""Generate document content from a note, then build downloadable
DOCX/PPTX/Markdown files.

Generation is a direct, tool-free LLM call (never routed through the agent
graph, which would defer ``organize_notes`` to a confirmation step instead of
producing content). Providers are resolved the same way the chatbot does:
the user's active ``LLMConfiguration`` first, then the ``NIM_*`` env fallback.
"""

import logging

from io import BytesIO

from django.conf import settings

from docx import Document
from docx.shared import Inches, Pt, RGBColor as DocxRGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as PptxInches

logger = logging.getLogger(__name__)

GENERATION_SYSTEM_PROMPT = (
    "You are a professional content writer. Produce the requested document "
    "using ONLY the supplied note as source material. "
    "Format rules:\n"
    "- Use `#` for the title, `##` for section headings, and `-` for bullets.\n"
    "- Presentations: keep each slide as a `##` section with 3-6 short bullets.\n"
    "- Articles: use an engaging intro, 2-4 `##` sections, and a short conclusion.\n"
    "- Output ONLY the document content. Never add preamble, commentary, "
    "questions, or a closing confirmation. Never ask the user to confirm or "
    "choose anything. Do not wrap the output in markdown code fences."
)

OPENAI_COMPATIBLE_PROVIDERS = {
    "openai",
    "nvidia",
    "huggingface",
    "custom",
    "groq",
    "openrouter",
    "mistral",
    "gemini",
    "deepseek",
    "together",
    "lmstudio",
}
KEYLESS_PROVIDERS = {"ollama", "lmstudio"}

PROVIDER_DEFAULT_ENDPOINTS = {
    "openai": "https://api.openai.com/v1",
    "nvidia": "https://integrate.api.nvidia.com/v1",
    "groq": "https://api.groq.com/openai/v1",
    "openrouter": "https://openrouter.ai/api/v1",
    "mistral": "https://api.mistral.ai/v1",
    "gemini": "https://generativelanguage.googleapis.com/v1beta/openai",
    "deepseek": "https://api.deepseek.com/v1",
    "glm": "https://open.bigmodel.cn/api/paas/v4",
    "together": "https://api.together.xyz/v1",
    "ollama": "http://localhost:11434",
    "lmstudio": "http://localhost:1234/v1",
}


def _resolve_llm(user) -> tuple[str, str, str, str]:
    """Return ``(api_key, base_url, model, provider)`` for the user's LLM.

    Prefers the user's active ``LLMConfiguration``; falls back to ``NIM_*``
    env settings so a demo deploy without a saved config still works.
    """
    from apps.integrations.models import LLMConfiguration

    config = (
        LLMConfiguration.objects.filter(user=user, is_active=True)
        .order_by("-updated_at")
        .first()
    )
    if config:
        endpoint = (config.api_endpoint or "").rstrip("/") or PROVIDER_DEFAULT_ENDPOINTS.get(
            config.provider, ""
        )
        return (
            config.decrypted_api_key,
            endpoint,
            config.model_name,
            config.provider,
        )
    return (
        getattr(settings, "NIM_API_KEY", ""),
        getattr(settings, "NIM_BASE_URL", "https://integrate.api.nvidia.com/v1"),
        getattr(settings, "NIM_MODEL", "meta/llama-3.1-8b-instruct"),
        "nvidia",
    )


def _call_llm(api_key: str, base_url: str, model: str, provider: str, prompt: str) -> str:
    """Single-shot tool-free LLM call returning the raw generated text."""
    if provider == "anthropic":
        from langchain_anthropic import ChatAnthropic

        llm = ChatAnthropic(
            model=model,
            api_key=api_key or "not-needed",
            temperature=0.7,
            max_tokens=4096,
            timeout=90,
            max_retries=1,
        )
    else:
        from langchain_openai import ChatOpenAI

        llm = ChatOpenAI(
            model=model,
            api_key=api_key or ("local-not-needed" if provider in KEYLESS_PROVIDERS else ""),
            base_url=base_url or None,
            temperature=0.7,
            max_tokens=4096,
            timeout=90,
            max_retries=1,
        )
    try:
        response = llm.invoke([{"role": "system", "content": GENERATION_SYSTEM_PROMPT}, {"role": "user", "content": prompt}])
    except Exception as exc:
        logger.warning(f"LLM call failed for provider={provider}, model={model}, url={base_url}: {exc}")
        msg = str(exc)
        if "401" in msg or "Unauthorized" in msg or "Authentication failed" in msg:
            raise RuntimeError(
                "The AI provider rejected your API key (401 Unauthorized). "
                "Check that the provider, model, endpoint and API key match. "
                "If you are using GLM, choose the 'Zhipu AI (GLM)' provider, not NVIDIA."
            ) from exc
        if "404" in msg:
            raise RuntimeError(
                f"The model '{model}' was not found at {base_url or 'the default endpoint'}. "
                "Verify the model name and endpoint for your provider."
            ) from exc
        raise RuntimeError(
            f"The AI request failed: {msg[:200]}. Check your LLM configuration and try again."
        ) from exc
    content = (getattr(response, "content", "") or "").strip()
    return content


def generate_content(
    user,
    note_title: str,
    note_content: str,
    target_format: str,
    style: str,
    max_length: int | None = None,
) -> str:
    """Generate full document content (markdown) for the given note."""
    style_desc = {
        "professional": "professional, clear and well structured",
        "creative": "creative and engaging",
        "academic": "academic with rigorous argumentation",
        "simple": "simple and easy to read",
    }.get(style, "professional, clear and well structured")

    prompt = (
        f"Write a {target_format} titled '{note_title}' "
        f"in a {style_desc} style, based on the note below.\n\n"
        f"--- NOTE ---\n{note_content}\n--- END NOTE ---\n"
    )
    if max_length:
        prompt += f"\nKeep the total output under {max_length} words."

    api_key, base_url, model, provider = _resolve_llm(user)
    if not api_key and provider not in KEYLESS_PROVIDERS:
        logger.warning("No LLM API key found for document generation")
        raise RuntimeError(
            "No AI model is configured. Add one under Settings → LLM Configuration."
        )
    content = _call_llm(api_key, base_url, model, provider, prompt)
    if not content:
        raise RuntimeError("The AI returned an empty response. Please try again.")
    return content

ACCENT = RGBColor(0x1F, 0x6F, 0xEB)
DARK = RGBColor(0x1F, 0x29, 0x37)


def _parse_blocks(content: str):
    """Split generated content into (kind, text) blocks.

    Recognizes markdown-ish headings and bullets produced by the agent.
    """
    blocks = []
    for raw_line in content.splitlines():
        line = raw_line.strip()
        if not line:
            continue
        if line.startswith("### "):
            blocks.append(("subhead", line[4:].strip()))
        elif line.startswith("## "):
            blocks.append(("head", line[3:].strip()))
        elif line.startswith("# "):
            blocks.append(("title", line[2:].strip()))
        elif line.startswith("- ") or line.startswith("* "):
            blocks.append(("bullet", line[2:].strip()))
        else:
            blocks.append(("para", line))
    return blocks


def _group_slides(blocks):
    """Group blocks into slides. A new `head` starts a new slide."""
    slides = []
    current = None
    for kind, text in blocks:
        if kind == "title" and current is None:
            current = {"title": text, "items": []}
        elif kind == "head":
            if current:
                slides.append(current)
            current = {"title": text, "items": []}
        elif current is None:
            current = {"title": "", "items": []}
        if kind in ("subhead", "bullet"):
            current["items"].append((kind, text))
        elif kind == "para":
            current["items"].append(("para", text))
    if current:
        slides.append(current)
    return slides


def _build_pptx(title: str, content: str) -> bytes:
    prs = Presentation()
    prs.slide_width = PptxInches(13.33)
    prs.slide_height = PptxInches(7.5)

    blank = prs.slide_layouts[6]

    slides = _group_slides(_parse_blocks(content))
    if not slides:
        slides = [{"title": title, "items": [("para", content[:400])]}]

    def add_title_slide():
        slide = prs.slides.add_slide(blank)
        tb = slide.shapes.add_textbox(PptxInches(1), PptxInches(2.2), PptxInches(11.3), PptxInches(2.5))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.alignment = PP_ALIGN.CENTER
        for run in p.runs:
            run.font.size = Pt(40)
            run.font.bold = True
            run.font.color.rgb = DARK
        p2 = tf.add_paragraph()
        p2.text = "Generated with ML-Auditor"
        p2.alignment = PP_ALIGN.CENTER
        for run in p2.runs:
            run.font.size = Pt(16)
            run.font.color.rgb = RGBColor(0x6B, 0x72, 0x80)

    add_title_slide()

    for slide_data in slides:
        slide = prs.slides.add_slide(blank)
        tb = slide.shapes.add_textbox(PptxInches(0.8), PptxInches(0.5), PptxInches(11.7), PptxInches(1.0))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = slide_data["title"]
        for run in p.runs:
            run.font.size = Pt(30)
            run.font.bold = True
            run.font.color.rgb = ACCENT

        body = slide.shapes.add_textbox(PptxInches(0.8), PptxInches(1.7), PptxInches(11.7), PptxInches(5.2))
        btf = body.text_frame
        btf.word_wrap = True
        first = True
        for kind, text in slide_data["items"]:
            if first:
                p = btf.paragraphs[0]
                first = False
            else:
                p = btf.add_paragraph()
            if kind == "subhead":
                p.text = text
                for run in p.runs:
                    run.font.bold = True
                    run.font.size = Pt(18)
                    run.font.color.rgb = DARK
            elif kind == "bullet":
                p.text = f"•  {text}"
                p.level = 0
                for run in p.runs:
                    run.font.size = Pt(18)
                    run.font.color.rgb = DARK
            else:
                p.text = text
                for run in p.runs:
                    run.font.size = Pt(18)
                    run.font.color.rgb = DARK

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _build_docx(title: str, content: str) -> bytes:
    doc = Document()

    h = doc.add_heading(title, level=0)
    for run in h.runs:
        run.font.color.rgb = DocxRGBColor(0x1F, 0x6F, 0xEB)

    for kind, text in _parse_blocks(content):
        if kind == "title":
            doc.add_heading(text, level=1)
        elif kind == "head":
            doc.add_heading(text, level=2)
        elif kind == "subhead":
            doc.add_heading(text, level=3)
        elif kind == "bullet":
            doc.add_paragraph(text, style="List Bullet")
        else:
            p = doc.add_paragraph(text)
            p.paragraph_format.space_after = Pt(6)

    buf = BytesIO()
    doc.save(buf)
    return buf.getvalue()


def build_document_bytes(title: str, content: str, file_format: str) -> bytes:
    """Return raw file bytes for the requested file_format (docx/pptx/md)."""
    if file_format == "docx":
        return _build_docx(title, content)
    if file_format == "pptx":
        return _build_pptx(title, content)
    return content.encode("utf-8")


def default_filename(title: str, file_format: str) -> str:
    safe = "".join(c if c.isalnum() or c in " -_" else "-" for c in title).strip().replace(" ", "_")
    return f"{safe or 'document'}.{file_format}"
